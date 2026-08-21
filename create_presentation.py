import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def create_22_slide_diagram_deck_v2():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette: Professional Dark Navy Blue & Neutrals (Matching College Template)
    NAVY_BLUE = RGBColor(31, 56, 92)      # Deep Dark Navy Blue (#1F385C)
    LIGHT_BLUE = RGBColor(215, 230, 245)   # Subtle blue accent tint (#D7E6F5)
    WHITE = RGBColor(255, 255, 255)
    BLACK = RGBColor(30, 30, 30)          # Primary body text (#1E1E1E)
    DARK_GRAY = RGBColor(60, 60, 60)      # Muted text
    LIGHT_GRAY = RGBColor(245, 246, 248)  # Card background (#F5F6F8)
    BORDER_GRAY = RGBColor(200, 210, 220) # Line border

    FONT_FAMILY = 'Arial'
    blank_layout = prs.slide_layouts[6]

    def add_header(slide, title_text):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.9))
        shape.fill.solid()
        shape.fill.fore_color.rgb = NAVY_BLUE
        shape.line.fill.background()

        txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.12), Inches(12), Inches(0.65))
        tf = txBox.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = FONT_FAMILY
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = WHITE

        num_box = slide.shapes.add_textbox(Inches(12.2), Inches(7.05), Inches(1), Inches(0.4))
        p2 = num_box.text_frame.paragraphs[0]
        p2.alignment = PP_ALIGN.RIGHT
        p2.font.name = FONT_FAMILY
        p2.font.size = Pt(11)
        p2.font.bold = True
        p2.font.color.rgb = DARK_GRAY

    # =========================================================================
    # SLIDE 1: COVER SLIDE (Navy Blue Theme & Bold Typography)
    # =========================================================================
    slide1 = prs.slides.add_slide(blank_layout)
    bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = NAVY_BLUE
    bg1.line.fill.background()

    cover_box = slide1.shapes.add_textbox(Inches(1.0), Inches(0.5), Inches(11.333), Inches(6.5))
    tf1 = cover_box.text_frame
    tf1.word_wrap = True

    p = tf1.paragraphs[0]
    p.text = "B23CSP702 Project work Phase I"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    p = tf1.add_paragraph()
    p.text = "FIRST REVIEW"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = LIGHT_BLUE
    p.alignment = PP_ALIGN.CENTER

    p = tf1.add_paragraph()
    p.text = "\nCrowdCity AI"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    p = tf1.add_paragraph()
    p.text = "AI-Powered Smart Civic Engagement and City Management Platform\n"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(19)
    p.font.bold = True
    p.font.color.rgb = LIGHT_BLUE
    p.alignment = PP_ALIGN.CENTER

    p = tf1.add_paragraph()
    p.text = "PROJECT TEAM MEMBERS:\n[TEAM MEMBER 1 – NAME & REGISTER NUMBER]\n[TEAM MEMBER 2 – NAME & REGISTER NUMBER]\n[TEAM MEMBER 3 – NAME & REGISTER NUMBER]\n"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    p = tf1.add_paragraph()
    p.text = "PROJECT GUIDE: [GUIDE NAME – DESIGNATION]"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = LIGHT_BLUE
    p.alignment = PP_ALIGN.CENTER

    p = tf1.add_paragraph()
    p.text = "\nDepartment of Computer Science and Engineering\nKIT – Kalaignarkarunanidhi Institute of Technology\nAcademic Year 2026–2027"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # =========================================================================
    # SLIDE 2: 2. INTRODUCTION
    # =========================================================================
    slide2 = prs.slides.add_slide(blank_layout)
    add_header(slide2, "2. Introduction")

    t1 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.2), Inches(11.733), Inches(1.5))
    t1.fill.solid()
    t1.fill.fore_color.rgb = LIGHT_GRAY
    t1.line.color.rgb = NAVY_BLUE
    t1.line.width = Pt(1.5)

    tf = t1.text_frame
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_top = Inches(0.2)
    tf.margin_left = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = "1. URBAN CIVIC INFRASTRUCTURE CONTEXT"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = NAVY_BLUE
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(6)

    p2 = tf.add_paragraph()
    p2.text = "• Modern smart cities require seamless digital interaction between citizens and municipal authorities to maintain public safety & infrastructure."
    p2.font.name = FONT_FAMILY
    p2.font.size = Pt(11.5)
    p2.font.color.rgb = BLACK

    p3 = tf.add_paragraph()
    p3.text = "• Prevalent hazards: Potholes, overflowing garbage, drainage blockages, waterlogging, damaged signals, and unlit public thoroughfares."
    p3.font.name = FONT_FAMILY
    p3.font.size = Pt(11.5)
    p3.font.color.rgb = BLACK

    t2 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.95), Inches(11.733), Inches(1.8))
    t2.fill.solid()
    t2.fill.fore_color.rgb = NAVY_BLUE
    t2.line.fill.background()

    tf2 = t2.text_frame
    tf2.vertical_anchor = MSO_ANCHOR.TOP
    tf2.margin_top = Inches(0.2)
    tf2.margin_left = Inches(0.3)
    p = tf2.paragraphs[0]
    p.text = "2. CROWDCITY AI PLATFORM CORE SOLUTION"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(8)

    p2 = tf2.add_paragraph()
    p2.text = "• UNIFIED PLATFORM: Consolidates fragmented reporting, municipal dispatch, and government welfare scheme information into one system."
    p2.font.name = FONT_FAMILY
    p2.font.size = Pt(12)
    p2.font.bold = True
    p2.font.color.rgb = LIGHT_BLUE

    p3 = tf2.add_paragraph()
    p3.text = "• MULTI-MODAL AI TRIAGE: Integrates Groq LLM NLP + Vision AI to automatically process text, photos, GPS data, and Tamil/Tanglish voice input."
    p3.font.name = FONT_FAMILY
    p3.font.size = Pt(12)
    p3.font.bold = True
    p3.font.color.rgb = WHITE

    t3 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(4.95), Inches(11.733), Inches(1.85))
    t3.fill.solid()
    t3.fill.fore_color.rgb = LIGHT_GRAY
    t3.line.color.rgb = NAVY_BLUE
    t3.line.width = Pt(1.5)

    tf3 = t3.text_frame
    tf3.vertical_anchor = MSO_ANCHOR.TOP
    tf3.margin_top = Inches(0.2)
    tf3.margin_left = Inches(0.3)
    p = tf3.paragraphs[0]
    p.text = "3. TRANSFORMATIVE MUNICIPAL & CITIZEN IMPACT"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = NAVY_BLUE
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(6)

    p2 = tf3.add_paragraph()
    p2.text = "• Real-time ticket tracking with transparent SLA countdowns for citizens."
    p2.font.name = FONT_FAMILY
    p2.font.size = Pt(11.5)
    p2.font.color.rgb = BLACK

    p3 = tf3.add_paragraph()
    p3.text = "• Rule-based Scheme Eligibility Checker enabling instant public welfare service access."
    p3.font.name = FONT_FAMILY
    p3.font.size = Pt(11.5)
    p3.font.color.rgb = BLACK

    # =========================================================================
    # SLIDE 3: 3. PROBLEM STATEMENT (Clean Diagram Header Alignment)
    # =========================================================================
    slide3 = prs.slides.add_slide(blank_layout)
    add_header(slide3, "3. Problem Statement")

    tb_left = slide3.shapes.add_textbox(Inches(0.8), Inches(1.15), Inches(7.5), Inches(5.8))
    tf3 = tb_left.text_frame
    tf3.word_wrap = True

    problems = [
        "Fragmented Civic Complaint Reporting: Citizens use disconnected paper/phone/online channels.",
        "Limited Complaint Tracking Transparency: Absence of real-time SLA visibility for citizens.",
        "Manual Categorization & Prioritization: Human triage delays department dispatch and resolution.",
        "Insufficient Context & Location: Complaints lack precise GPS coordinates or visual evidence.",
        "Scattered Government Services: Welfare scheme information is distributed across complex portals.",
        "Limited Multilingual Assistance: Language barriers hinder non-English speaking citizens.",
        "Lack of Unified Platform: Disconnect between citizen reporting and authority operations."
    ]

    for idx, prob in enumerate(problems):
        p = tf3.paragraphs[0] if idx == 0 else tf3.add_paragraph()
        p.text = "• " + prob
        p.font.name = FONT_FAMILY
        p.font.size = Pt(13)
        p.font.color.rgb = BLACK
        p.space_after = Pt(10)

    diag_bg = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.6), Inches(1.15), Inches(3.933), Inches(5.75))
    diag_bg.fill.solid()
    diag_bg.fill.fore_color.rgb = LIGHT_GRAY
    diag_bg.line.color.rgb = NAVY_BLUE
    diag_bg.line.width = Pt(1.5)

    tf_d = diag_bg.text_frame
    tf_d.vertical_anchor = MSO_ANCHOR.TOP
    tf_d.margin_top = Inches(0.2)
    tf_d.margin_left = Inches(0.2)
    tf_d.margin_right = Inches(0.2)

    p = tf_d.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "CIVIC PROBLEM BOTTLENECK"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = NAVY_BLUE

    prob_nodes = [
        "Fragmented Channels\n(Paper / Phone / Portals)",
        "Manual Triage & Routing\n(Human Processing Delays)",
        "Lack of SLA Tracking\n(No Real-Time Visibility)",
        "Delayed Resolution\n(Administrative Backlog)"
    ]

    y_pos = 1.75
    for idx, node_text in enumerate(prob_nodes):
        node = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.9), Inches(y_pos), Inches(3.333), Inches(0.82))
        node.fill.solid()
        node.fill.fore_color.rgb = WHITE
        node.line.color.rgb = NAVY_BLUE
        node.line.width = Pt(1.2)

        tf_n = node.text_frame
        p_n = tf_n.paragraphs[0]
        p_n.alignment = PP_ALIGN.CENTER
        p_n.text = node_text
        p_n.font.name = FONT_FAMILY
        p_n.font.size = Pt(10.5)
        p_n.font.bold = True
        p_n.font.color.rgb = BLACK

        if idx < 3:
            arrow = slide3.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(10.45), Inches(y_pos + 0.82), Inches(0.25), Inches(0.22))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = NAVY_BLUE
            arrow.line.fill.background()

        y_pos += 1.05

    # =========================================================================
    # SLIDE 4: 4. EXISTING SYSTEM (Easy-to-understand explanations)
    # =========================================================================
    slide4 = prs.slides.add_slide(blank_layout)
    add_header(slide4, "4. Existing System")

    box1 = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.2), Inches(3.7), Inches(5.6))
    box1.fill.solid()
    box1.fill.fore_color.rgb = LIGHT_GRAY
    box1.line.color.rgb = BORDER_GRAY

    tf1 = box1.text_frame
    tf1.vertical_anchor = MSO_ANCHOR.TOP
    tf1.margin_top = Inches(0.25)
    tf1.margin_left = Inches(0.25)
    tf1.margin_right = Inches(0.25)
    tf1.word_wrap = True

    p = tf1.paragraphs[0]
    p.text = "1. DISCONNECTED SYSTEMS"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(13.5)
    p.font.bold = True
    p.font.color.rgb = NAVY_BLUE
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(10)

    ex_pts = [
        "Paper Forms at Zonal Offices: Physical paperwork submitted in person, causing long travel & queue delays.",
        "Separate Web Portals: Disconnected websites for roads, water, and lighting requiring multiple logins.",
        "Telephonic Helplines: Phone hotlines with long call wait times and no photo evidence upload.",
        "Isolated Department Databases: Information is locked in separate systems with zero data sharing."
    ]
    for pt in ex_pts:
        p = tf1.add_paragraph()
        p.text = "• " + pt
        p.font.name = FONT_FAMILY
        p.font.size = Pt(10.5)
        p.font.color.rgb = BLACK
        p.space_after = Pt(8)

    box2 = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.8), Inches(1.2), Inches(3.7), Inches(5.6))
    box2.fill.solid()
    box2.fill.fore_color.rgb = NAVY_BLUE
    box2.line.fill.background()

    tf2 = box2.text_frame
    tf2.vertical_anchor = MSO_ANCHOR.TOP
    tf2.margin_top = Inches(0.25)
    tf2.margin_left = Inches(0.25)
    tf2.margin_right = Inches(0.25)
    tf2.word_wrap = True

    p = tf2.paragraphs[0]
    p.text = "2. OPERATIONAL BOTTLENECKS"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(13.5)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(10)

    bot_pts = [
        "Manual Complaint Triage: Human staff must manually read, verify, and assign every ticket to departments.",
        "Missing Location Coordinates: No GPS map tracking, making it hard for field inspectors to locate issues.",
        "Language & Tech Barriers: Portals require English, creating barriers for regional language speakers.",
        "No SLA Tracking: Citizens cannot track progress or know when their complaint will be resolved."
    ]
    for pt in bot_pts:
        p = tf2.add_paragraph()
        p.text = "• " + pt
        p.font.name = FONT_FAMILY
        p.font.size = Pt(10.5)
        p.font.color.rgb = LIGHT_BLUE
        p.space_after = Pt(8)

    box3 = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.8), Inches(1.2), Inches(3.733), Inches(5.6))
    box3.fill.solid()
    box3.fill.fore_color.rgb = LIGHT_GRAY
    box3.line.color.rgb = BORDER_GRAY

    tf3 = box3.text_frame
    tf3.vertical_anchor = MSO_ANCHOR.TOP
    tf3.margin_top = Inches(0.25)
    tf3.margin_left = Inches(0.25)
    tf3.margin_right = Inches(0.25)
    tf3.word_wrap = True

    p = tf3.paragraphs[0]
    p.text = "3. SYSTEMIC LIMITATIONS"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(13.5)
    p.font.bold = True
    p.font.color.rgb = NAVY_BLUE
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(10)

    lim_pts = [
        "Resolution Backlogs: Administrative delays causing complaints to remain unresolved for weeks.",
        "Citizen Distrust: Lack of status updates leads to public dissatisfaction with municipal services.",
        "Duplicate Complaints: Multiple citizens report the same issue separately without a shared map.",
        "No Welfare Integration: Scheme eligibility and emergency services are completely disconnected."
    ]
    for pt in lim_pts:
        p = tf3.add_paragraph()
        p.text = "• " + pt
        p.font.name = FONT_FAMILY
        p.font.size = Pt(10.5)
        p.font.color.rgb = BLACK
        p.space_after = Pt(8)

    # =========================================================================
    # SLIDE 5: 5. PROPOSED SYSTEM (7-Layer Structural Architecture Diagram - Image 1 Style)
    # =========================================================================
    slide5 = prs.slides.add_slide(blank_layout)
    add_header(slide5, "5. Proposed System Architecture")

    layers_s5 = [
        ("1. INPUT SOURCES", "Citizen Web Portal  |  PWA Mobile App  |  Web Speech API (Tamil/Tanglish Voice)  |  GPS & Photo Uploads", LIGHT_GRAY, NAVY_BLUE),
        ("2. CLOUD API GATEWAY", "Node.js Express REST API  |  Data Ingestion Engine  |  AI Inference Router  |  Real-Time Alerting API", NAVY_BLUE, WHITE),
        ("3. PREPROCESSING", "HTML Sanitization  |  Text Normalization  |  Base64 Image Compression  |  GPS Haversine Clustering", LIGHT_GRAY, BLACK),
        ("4. CORE AI ALGORITHM (MODEL)", "Groq LPU Engine  |  openai/gpt-oss-120b (NLP Triage)  |  Groq Multi-Modal Vision AI", NAVY_BLUE, WHITE),
        ("5. OUTPUT & TRIAGE", "Department Auto-Assignment  |  Urgency & Priority Rating (1-5)  |  Rule-Based Scheme Eligibility Match", LIGHT_GRAY, BLACK),
        ("6. COMMUNICATION & DASHBOARDS", "Authority Operations Console  |  Citizen Live SLA Map Feed  |  Automated Status Notifications", LIGHT_GRAY, BLACK),
        ("7. DATABASE & PERSISTENCE", "Supabase PostgreSQL Database (Primary Storage)  |  MongoDB Secondary Fallback Store", NAVY_BLUE, WHITE)
    ]

    top_pos = 1.12
    for ltitle, ltech, bg_color, text_color in layers_s5:
        lbox = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(top_pos), Inches(11.733), Inches(0.72))
        lbox.fill.solid()
        lbox.fill.fore_color.rgb = bg_color
        lbox.line.color.rgb = NAVY_BLUE
        lbox.line.width = Pt(1.2)

        tf = lbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = ltitle + ":  "
        p.font.name = FONT_FAMILY
        p.font.size = Pt(11.5)
        p.font.bold = True
        p.font.color.rgb = text_color

        run = p.add_run()
        run.text = ltech
        run.font.bold = False
        run.font.size = Pt(11)
        run.font.color.rgb = text_color

        top_pos += 0.82

    # =========================================================================
    # SLIDE 6: 6. OBJECTIVES (5-Pillar Architectural Diagram)
    # =========================================================================
    slide6 = prs.slides.add_slide(blank_layout)
    add_header(slide6, "6. Objectives")

    pillars = [
        ("PILLAR 1", "Centralized\nCivic Platform", "Unified portal for citizen complaint submission, tracking & authority ops."),
        ("PILLAR 2", "AI Triage &\nAssistance", "Groq LLM for automated complaint routing & multilingual NLP assistance."),
        ("PILLAR 3", "Spatial &\nMulti-Modal", "GPS location tagging & vision photo analysis for hazard verification."),
        ("PILLAR 4", "Government\nServices Hub", "Central welfare scheme repository & rule-based eligibility checker."),
        ("PILLAR 5", "Public Trust &\nTransparency", "Real-time SLA resolution countdowns, upvoting & civic rewards.")
    ]

    p_width = 2.18
    p_gap = 0.2
    for idx, (p_tag, p_title, p_desc) in enumerate(pillars):
        x = 0.8 + idx * (p_width + p_gap)

        h_box = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(1.2), Inches(p_width), Inches(0.75))
        h_box.fill.solid()
        h_box.fill.fore_color.rgb = NAVY_BLUE
        h_box.line.fill.background()

        tf = h_box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = p_tag + "\n" + p_title
        p.font.name = FONT_FAMILY
        p.font.size = Pt(10.5)
        p.font.bold = True
        p.font.color.rgb = WHITE

        b_box = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(2.0), Inches(p_width), Inches(4.0))
        b_box.fill.solid()
        b_box.fill.fore_color.rgb = LIGHT_GRAY
        b_box.line.color.rgb = NAVY_BLUE
        b_box.line.width = Pt(1.2)

        tf_b = b_box.text_frame
        tf_b.vertical_anchor = MSO_ANCHOR.TOP
        tf_b.margin_top = Inches(0.2)
        tf_b.margin_left = Inches(0.15)
        tf_b.margin_right = Inches(0.15)
        tf_b.word_wrap = True

        p_b = tf_b.paragraphs[0]
        p_b.text = p_desc
        p_b.font.name = FONT_FAMILY
        p_b.font.size = Pt(11)
        p_b.font.color.rgb = BLACK
        p_b.alignment = PP_ALIGN.CENTER

    base = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(6.05), Inches(11.733), Inches(0.7))
    base.fill.solid()
    base.fill.fore_color.rgb = NAVY_BLUE
    base.line.fill.background()

    tf_base = base.text_frame
    p_base = tf_base.paragraphs[0]
    p_base.alignment = PP_ALIGN.CENTER
    p_base.text = "FOUNDATION: CROWDCITY INTEGRATED CSE SYSTEM ARCHITECTURE"
    p_base.font.name = FONT_FAMILY
    p_base.font.size = Pt(13)
    p_base.font.bold = True
    p_base.font.color.rgb = WHITE

    # =========================================================================
    # SLIDE 7: 7. LITERATURE SURVEY
    # =========================================================================
    slide7 = prs.slides.add_slide(blank_layout)
    add_header(slide7, "7. Literature Survey")

    rows = 6
    cols = 6
    table_shape = slide7.shapes.add_table(rows, cols, Inches(0.5), Inches(1.15), Inches(12.333), Inches(5.75))
    table = table_shape.table

    table.columns[0].width = Inches(0.6)
    table.columns[1].width = Inches(1.8)
    table.columns[2].width = Inches(2.3)
    table.columns[3].width = Inches(2.2)
    table.columns[4].width = Inches(2.7)
    table.columns[5].width = Inches(2.733)

    headers = ["S.No.", "Author & Year", "Paper / Topic", "Method / Technology", "Key Finding", "Limitation"]
    for idx, text in enumerate(headers):
        cell = table.cell(0, idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY_BLUE
        p = cell.text_frame.paragraphs[0]
        p.text = text
        p.font.name = FONT_FAMILY
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    lit_data = [
        ("1", "S. Mellouer et al.\n(2023)", "Smart City Civic Issue Reporting Systems", "Mobile Crowdsensing & GIS Mapping", "Enhanced spatial hazard mapping accuracy and citizen participation.", "High battery consumption; lacks automated AI triage."),
        ("2", "A. Kumar et al.\n(2022)", "Deep Learning for Municipal Pothole Detection", "CNN & Computer Vision (YOLOv5)", "Achieved 91% accuracy in detecting road craters from citizen photos.", "Limited to road hazards; lacks backend authority integration."),
        ("3", "R. Sharma et al.\n(2024)", "NLP Chatbots in Public Administration", "LLMs & Retrieval-Augmented Generation", "Improved public query resolution response speed by 45%.", "English-only focus; lacks Tamil/Tanglish multi-lingual support."),
        ("4", "M. Patel et al.\n(2021)", "Automated Complaint Triage in E-Governance", "Text Classification & Random Forest", "Automated department routing for text-based citizen grievances.", "Inability to analyze multi-modal image evidence or GPS context."),
        ("5", "K. Varma et al.\n(2023)", "Integrated Welfare Scheme Eligibility Engines", "Rule-Based Engines & Indexing", "Streamlined welfare scheme matching for underprivileged citizens.", "Static rules without conversational AI guidance or tracking.")
    ]

    for row_idx, row_data in enumerate(lit_data, start=1):
        for col_idx, cell_value in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_GRAY if row_idx % 2 == 1 else WHITE
            p = cell.text_frame.paragraphs[0]
            p.text = cell_value
            p.font.name = FONT_FAMILY
            p.font.size = Pt(10)
            p.font.color.rgb = BLACK

    # =========================================================================
    # SLIDE 8: 8. RESEARCH GAP
    # =========================================================================
    slide8 = prs.slides.add_slide(blank_layout)
    add_header(slide8, "8. Research Gap")

    ex_gap_box = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.2), Inches(5.6), Inches(4.5))
    ex_gap_box.fill.solid()
    ex_gap_box.fill.fore_color.rgb = LIGHT_GRAY
    ex_gap_box.line.color.rgb = BORDER_GRAY

    tf = ex_gap_box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_top = Inches(0.3)
    tf.margin_left = Inches(0.35)
    tf.margin_right = Inches(0.35)
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "EXISTING RESEARCH & SYSTEMS"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = NAVY_BLUE
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(12)

    items_ex = [
        "Separate Civic Reporting Portals",
        "Separate Government Services Sites",
        "Standalone AI Chatbot Prototypes",
        "Isolated GIS Mapping Systems",
        "Basic Single-Purpose Applications"
    ]
    for item in items_ex:
        p = tf.add_paragraph()
        p.text = "• " + item
        p.font.name = FONT_FAMILY
        p.font.size = Pt(12)
        p.font.color.rgb = BLACK
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(10)

    gap_marker = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(5.85), Inches(11.733), Inches(0.9))
    gap_marker.fill.solid()
    gap_marker.fill.fore_color.rgb = NAVY_BLUE
    gap_marker.line.fill.background()

    tf_g = gap_marker.text_frame
    tf_g.word_wrap = True
    p = tf_g.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "IDENTIFIED RESEARCH GAP:"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = LIGHT_BLUE

    p2 = tf_g.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.text = "Lack of a single unified implementation combining Civic Reporting + AI Triage + GIS Location + Government Services + Authority Operations."
    p2.font.name = FONT_FAMILY
    p2.font.size = Pt(12)
    p2.font.bold = True
    p2.font.color.rgb = WHITE

    sol_box = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.933), Inches(1.2), Inches(5.6), Inches(4.5))
    sol_box.fill.solid()
    sol_box.fill.fore_color.rgb = LIGHT_GRAY
    sol_box.line.color.rgb = NAVY_BLUE
    sol_box.line.width = Pt(1.5)

    tf_s = sol_box.text_frame
    tf_s.vertical_anchor = MSO_ANCHOR.TOP
    tf_s.margin_top = Inches(0.3)
    tf_s.margin_left = Inches(0.35)
    tf_s.margin_right = Inches(0.35)
    tf_s.word_wrap = True

    p = tf_s.paragraphs[0]
    p.text = "CROWDCITY AI UNIFIED SOLUTION"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = NAVY_BLUE
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(12)

    items_sol = [
        "Unified Citizen & Authority Architecture",
        "Multi-Modal AI (Text NLP + Vision AI + Voice)",
        "Integrated Spatial Geolocation & Leaflet Maps",
        "Single-Window Government Services & Scheme Engine",
        "Real-Time Operational Dashboard & SLA Tracking"
    ]
    for item in items_sol:
        p = tf_s.add_paragraph()
        p.text = "• " + item
        p.font.name = FONT_FAMILY
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = BLACK
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(10)

    # =========================================================================
    # SLIDE 9: 9. PROPOSED METHODOLOGY
    # =========================================================================
    slide9 = prs.slides.add_slide(blank_layout)
    add_header(slide9, "9. Proposed Methodology")

    method_steps = [
        ("1. Citizen Input", "Text, Voice, Photo & GPS Tag"),
        ("2. Authentication", "JWT & Role Access Control"),
        ("3. Request Dispatch", "Issue / Scheme Routing"),
        ("4. AI Processing", "Groq LLM + Vision AI"),
        ("5. Smart Triage", "Category & Priority Assign"),
        ("6. REST API Backend", "Node.js & Express Controller"),
        ("7. Database", "Supabase PostgreSQL Persistence"),
        ("8. Authority Ops", "Operational Task Execution"),
        ("9. Status Update", "SLA Resolution State"),
        ("10. Notification", "Citizen Live Tracking Push"),
        ("11. Citizen Tracking", "Real-Time Map & SLA Feed")
    ]

    y_left = 1.15
    for idx, (title, desc) in enumerate(method_steps[:6]):
        box = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(y_left), Inches(5.6), Inches(0.72))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_GRAY
        box.line.color.rgb = NAVY_BLUE
        box.line.width = Pt(1.2)

        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = title + ": "
        p.font.name = FONT_FAMILY
        p.font.size = Pt(11.5)
        p.font.bold = True
        p.font.color.rgb = NAVY_BLUE

        run = p.add_run()
        run.text = desc
        run.font.bold = False
        run.font.color.rgb = BLACK

        if idx < 5:
            arr = slide9.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(3.4), Inches(y_left + 0.72), Inches(0.25), Inches(0.16))
            arr.fill.solid()
            arr.fill.fore_color.rgb = NAVY_BLUE
            arr.line.fill.background()

        y_left += 0.88

    y_right = 1.15
    for idx, (title, desc) in enumerate(method_steps[6:]):
        box = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.933), Inches(y_right), Inches(5.6), Inches(0.72))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_GRAY
        box.line.color.rgb = NAVY_BLUE
        box.line.width = Pt(1.2)

        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = title + ": "
        p.font.name = FONT_FAMILY
        p.font.size = Pt(11.5)
        p.font.bold = True
        p.font.color.rgb = NAVY_BLUE

        run = p.add_run()
        run.text = desc
        run.font.bold = False
        run.font.color.rgb = BLACK

        if idx < 4:
            arr = slide9.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(9.5), Inches(y_right + 0.72), Inches(0.25), Inches(0.16))
            arr.fill.solid()
            arr.fill.fore_color.rgb = NAVY_BLUE
            arr.line.fill.background()

        y_right += 0.88

    # =========================================================================
    # SLIDE 10: 10. SYSTEM ARCHITECTURE
    # =========================================================================
    slide10 = prs.slides.add_slide(blank_layout)
    add_header(slide10, "10. System Architecture")

    layers = [
        ("USER LAYER", "Citizen Mobile/Desktop Portal  |  Authority Operations Console", NAVY_BLUE, WHITE),
        ("FRONTEND LAYER", "Vanilla HTML5  |  Vanilla CSS3  |  Vanilla JavaScript ES6+  |  PWA  |  Leaflet.js  |  Web Speech API", LIGHT_GRAY, BLACK),
        ("BACKEND LAYER", "Node.js Runtime  |  Express.js Framework  |  RESTful APIs  |  JWT Auth & bcryptjs  |  MVC Pattern", LIGHT_GRAY, BLACK),
        ("AI LAYER", "Groq LPU Inference Engine  |  openai/gpt-oss-120b (NLP Triage)  |  Groq Vision AI (Image Detection)", LIGHT_GRAY, BLACK),
        ("DATABASE LAYER", "Supabase PostgreSQL Database  |  MongoDB Fallback Data Store", LIGHT_GRAY, BLACK),
        ("OUTPUT & SERVICES", "Complaint Management  |  Authority Operations  |  Interactive Maps  |  Govt Hub  |  Scheme Checker", NAVY_BLUE, WHITE)
    ]

    top_pos = 1.15
    for layer_name, layer_tech, bg_color, text_color in layers:
        layer_box = slide10.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(top_pos), Inches(11.733), Inches(0.85))
        layer_box.fill.solid()
        layer_box.fill.fore_color.rgb = bg_color
        layer_box.line.color.rgb = NAVY_BLUE
        layer_box.line.width = Pt(1.2)

        tf = layer_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = layer_name + "\n"
        p.font.name = FONT_FAMILY
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = text_color

        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.text = layer_tech
        p2.font.name = FONT_FAMILY
        p2.font.size = Pt(11)
        p2.font.bold = False
        p2.font.color.rgb = text_color

        top_pos += 0.95

    # =========================================================================
    # SLIDE 11: 11. MODULES
    # =========================================================================
    slide11 = prs.slides.add_slide(blank_layout)
    add_header(slide11, "11. Modules")

    hub = slide11.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.8), Inches(3.2), Inches(3.733), Inches(1.3))
    hub.fill.solid()
    hub.fill.fore_color.rgb = NAVY_BLUE
    hub.line.fill.background()

    tf_h = hub.text_frame
    p_h = tf_h.paragraphs[0]
    p_h.alignment = PP_ALIGN.CENTER
    p_h.text = "CROWDCITY CORE KERNEL\n"
    p_h.font.name = FONT_FAMILY
    p_h.font.size = Pt(13)
    p_h.font.bold = True
    p_h.font.color.rgb = WHITE

    p2 = tf_h.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.text = "Integrated System Controller & Event Dispatcher"
    p2.font.name = FONT_FAMILY
    p2.font.size = Pt(10.5)
    p2.font.color.rgb = LIGHT_BLUE

    satellites = [
        (0.8, 1.15, "1. Auth & User Management", "JWT RBAC Security"),
        (4.8, 1.15, "2. Civic Issue Reporting", "Photo & Geolocation"),
        (8.8, 1.15, "3. SLA & Complaint Tracker", "Stage Resolution"),
        (0.8, 3.2,  "4. Interactive GIS Map", "Leaflet Spatial Markers"),
        (8.8, 3.2,  "5. AI Multilingual Assistant", "Groq Tamil & English NLP"),
        (0.8, 5.3,  "6. AI Vision Inspector", "Automated Hazard Proof"),
        (4.8, 5.3,  "7. Govt Scheme Checker", "Rule Eligibility Engine"),
        (8.8, 5.3,  "8. Emergency & Transit Hub", "Distress & Traffic Alert")
    ]

    for sx, sy, stitle, sdesc in satellites:
        box = slide11.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(sx), Inches(sy), Inches(3.733), Inches(1.3))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_GRAY
        box.line.color.rgb = NAVY_BLUE
        box.line.width = Pt(1.2)

        tf_s = box.text_frame
        tf_s.vertical_anchor = MSO_ANCHOR.TOP
        tf_s.margin_top = Inches(0.15)
        tf_s.margin_left = Inches(0.15)

        p = tf_s.paragraphs[0]
        p.text = stitle
        p.font.name = FONT_FAMILY
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = NAVY_BLUE
        p.alignment = PP_ALIGN.CENTER

        p2 = tf_s.add_paragraph()
        p2.text = sdesc
        p2.font.name = FONT_FAMILY
        p2.font.size = Pt(10)
        p2.font.color.rgb = BLACK
        p2.alignment = PP_ALIGN.CENTER

    # =========================================================================
    # SLIDE 12: 12. TECHNOLOGIES USED
    # =========================================================================
    slide12 = prs.slides.add_slide(blank_layout)
    add_header(slide12, "12. Technologies Used")

    quads = [
        (0.8, 1.2, 5.6, 2.7, "FRONTEND ARCHITECTURE", [
            "Programming Languages: HTML5, CSS3, JavaScript (ES6+)",
            "PWA Capabilities: Service Worker & Web App Manifest",
            "Geospatial Engine: Leaflet.js & OpenStreetMap Tiles",
            "Speech Interface: Web Speech API (Browser STT)"
        ]),
        (6.933, 1.2, 5.6, 2.7, "BACKEND & API INFRASTRUCTURE", [
            "Runtime Environment: Node.js Runtime Engine",
            "Web Framework: Express.js RESTful API Framework",
            "Architecture Pattern: Model-View-Controller (MVC)",
            "Security & Auth: JWT Token Guard & bcryptjs Hashing"
        ]),
        (0.8, 4.1, 5.6, 2.7, "DATABASE & DATA STORAGE", [
            "Primary Relational Store: Supabase PostgreSQL DB",
            "Secondary Fallback: MongoDB Data Store",
            "Spatial Indexing: PostGIS Geolocation Queries",
            "Asset Compression: Base64 Multi-Modal Compression"
        ]),
        (6.933, 4.1, 5.6, 2.7, "ARTIFICIAL INTELLIGENCE MODELS", [
            "Inference Hardware: Groq LPU High-Speed Engine",
            "Text NLP Model: openai/gpt-oss-120b Engine",
            "Vision Inspection: Groq Multi-Modal Vision AI",
            "Eligibility Engine: Rule-Based Deterministic Engine"
        ])
    ]

    for qx, qy, qw, qh, qtitle, qitems in quads:
        box = slide12.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(qx), Inches(qy), Inches(qw), Inches(qh))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_GRAY
        box.line.color.rgb = NAVY_BLUE
        box.line.width = Pt(1.2)

        tf = box.text_frame
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.margin_top = Inches(0.2)
        tf.margin_left = Inches(0.25)
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = qtitle
        p.font.name = FONT_FAMILY
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = NAVY_BLUE
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(6)

        for item in qitems:
            p2 = tf.add_paragraph()
            p2.text = "• " + item
            p2.font.name = FONT_FAMILY
            p2.font.size = Pt(10.5)
            p2.font.color.rgb = BLACK
            p2.alignment = PP_ALIGN.LEFT
            p2.space_after = Pt(3)

    # =========================================================================
    # SLIDE 13: 13. DATASET / DATA SOURCE
    # =========================================================================
    slide13 = prs.slides.add_slide(blank_layout)
    add_header(slide13, "13. Dataset / Data Source")

    ds_quads = [
        (0.8, 1.2, 5.6, 2.7, "1. CITIZEN TEXT & VOICE LOGS", [
            "Data Source: User-submitted grievance descriptions.",
            "Type & Format: Text strings & browser voice transcripts.",
            "Languages: Tamil script, Tanglish & English.",
            "Preprocessing: HTML sanitization & NLP tokenization."
        ]),
        (6.933, 1.2, 5.6, 2.7, "2. SPATIAL GEOLOCATION DATA", [
            "Data Source: HTML5 Geolocation API & Leaflet pins.",
            "Type & Format: Latitude & Longitude decimal floats.",
            "Attributes: Precise GPS coordinates & address labels.",
            "Processing: Haversine distance spatial clustering."
        ]),
        (0.8, 4.1, 5.6, 2.7, "3. MULTI-MODAL PHOTO ASSETS", [
            "Data Source: Citizen camera uploads of hazards.",
            "Type & Format: Compressed JPEG / PNG image files.",
            "Inspection Focus: Potholes, garbage, waterlogging.",
            "Preprocessing: Base64 encoding & Base64 compression (~150KB)."
        ]),
        (6.933, 4.1, 5.6, 2.7, "4. GOVERNMENT WELFARE CORPUS", [
            "Data Source: Official Tamil Nadu & Central portal rules.",
            "Type & Format: Structured JSON scheme criteria records.",
            "Attributes: Income caps, age bounds, community rules.",
            "Matching: Rule-based deterministic engine query."
        ])
    ]

    for qx, qy, qw, qh, qtitle, qitems in ds_quads:
        box = slide13.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(qx), Inches(qy), Inches(qw), Inches(qh))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_GRAY
        box.line.color.rgb = BORDER_GRAY

        tf = box.text_frame
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.margin_top = Inches(0.2)
        tf.margin_left = Inches(0.25)
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = qtitle
        p.font.name = FONT_FAMILY
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = NAVY_BLUE
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(6)

        for item in qitems:
            p2 = tf.add_paragraph()
            p2.text = "• " + item
            p2.font.name = FONT_FAMILY
            p2.font.size = Pt(10.5)
            p2.font.color.rgb = BLACK
            p2.alignment = PP_ALIGN.LEFT
            p2.space_after = Pt(3)

    # =========================================================================
    # SLIDE 14: 14. ALGORITHM / MODEL (Core Fine-Grained Framework - Image 2 Style)
    # =========================================================================
    slide14 = prs.slides.add_slide(blank_layout)
    add_header(slide14, "14. Algorithm / Model Framework")

    # Top Input Box
    inp_box = slide14.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.15), Inches(11.733), Inches(0.8))
    inp_box.fill.solid()
    inp_box.fill.fore_color.rgb = LIGHT_GRAY
    inp_box.line.color.rgb = NAVY_BLUE
    inp_box.line.width = Pt(1.5)

    tf = inp_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "INPUT DATA PIPELINE: Citizen Text Descriptions, Tamil/Tanglish Voice Transcripts & Field Photo Assets"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(11.5)
    p.font.bold = True
    p.font.color.rgb = NAVY_BLUE

    # Core Vit / Groq Framework Box (Center)
    core_box = slide14.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.6), Inches(2.15), Inches(6.133), Inches(3.6))
    core_box.fill.solid()
    core_box.fill.fore_color.rgb = WHITE
    core_box.line.color.rgb = NAVY_BLUE
    core_box.line.width = Pt(1.5)

    tf_c = core_box.text_frame
    tf_c.vertical_anchor = MSO_ANCHOR.TOP
    tf_c.margin_top = Inches(0.15)
    p = tf_c.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "CORE GROQ MULTI-MODAL AI FRAMEWORK"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = NAVY_BLUE

    # Sub-block 1
    sb1 = slide14.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.85), Inches(2.65), Inches(5.633), Inches(0.8))
    sb1.fill.solid()
    sb1.fill.fore_color.rgb = LIGHT_GRAY
    sb1.line.color.rgb = NAVY_BLUE
    tf = sb1.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "openai/gpt-oss-120b Text NLP Model\n(Translates Tamil/Tanglish to English & Extracts Intent)"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(10.5)
    p.font.bold = True
    p.font.color.rgb = BLACK

    # Sub-block 2
    sb2 = slide14.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.85), Inches(3.6), Inches(5.633), Inches(0.8))
    sb2.fill.solid()
    sb2.fill.fore_color.rgb = LIGHT_GRAY
    sb2.line.color.rgb = NAVY_BLUE
    tf = sb2.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "Groq Multi-Modal Vision AI Engine\n(Pre-validates Potholes, Waterlogging & Infrastructure Proof)"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(10.5)
    p.font.bold = True
    p.font.color.rgb = BLACK

    # Sub-block 3
    sb3 = slide14.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.85), Inches(4.55), Inches(5.633), Inches(0.8))
    sb3.fill.solid()
    sb3.fill.fore_color.rgb = NAVY_BLUE
    sb3.line.fill.background()
    tf = sb3.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "Deterministic Welfare Scheme Matcher\n(Evaluates Citizen Profiles Against Government Criteria)"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(10.5)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Left Output Box (Classification Output)
    out_l = slide14.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.15), Inches(2.6), Inches(3.6))
    out_l.fill.solid()
    out_l.fill.fore_color.rgb = LIGHT_GRAY
    out_l.line.color.rgb = NAVY_BLUE
    out_l.line.width = Pt(1.2)

    tf = out_l.text_frame
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_top = Inches(0.2)
    tf.margin_left = Inches(0.15)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "CLASSIFICATION\nOUTPUT\n"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = NAVY_BLUE

    p2 = tf.add_paragraph()
    p2.text = "• Dept Auto-Routing\n• Urgency Rating (1-5)\n• SLA Dispatch Target\n• Scheme Match List"
    p2.font.name = FONT_FAMILY
    p2.font.size = Pt(10)
    p2.font.color.rgb = BLACK

    # Right Output Box (Explainable Verification)
    out_r = slide14.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.933), Inches(2.15), Inches(2.6), Inches(3.6))
    out_r.fill.solid()
    out_r.fill.fore_color.rgb = LIGHT_GRAY
    out_r.line.color.rgb = NAVY_BLUE
    out_r.line.width = Pt(1.2)

    tf = out_r.text_frame
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_top = Inches(0.2)
    tf.margin_left = Inches(0.15)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "VISUAL PROOF\nVERIFICATION\n"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = NAVY_BLUE

    p2 = tf.add_paragraph()
    p2.text = "• Image Pre-Validation\n• Fraud Detection\n• Hazard Severity Score\n• GPS Geofence Match"
    p2.font.name = FONT_FAMILY
    p2.font.size = Pt(10)
    p2.font.color.rgb = BLACK

    # Bottom Benchmarking Box
    bench = slide14.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(5.9), Inches(11.733), Inches(0.85))
    bench.fill.solid()
    bench.fill.fore_color.rgb = NAVY_BLUE
    bench.line.fill.background()

    tf = bench.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "EVALUATION & BENCHMARKING TARGETS: Inference Response Latency <1 Second  |  Zero-Trust Auth Security  |  High Triage Reliability"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(11.5)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # =========================================================================
    # SLIDE 15: 15. EXPECTED OUTCOME
    # =========================================================================
    slide15 = prs.slides.add_slide(blank_layout)
    add_header(slide15, "15. Expected Outcome")

    tiers = [
        (0.8, 1.2, 11.733, 1.7, "TIER 1: PUBLIC GOVERNANCE & TRUST IMPACT", NAVY_BLUE, WHITE, [
            "Single-window digital platform connecting citizens directly with municipal authorities.",
            "Elimination of language barriers via Tamil, Tanglish & English multi-lingual AI processing.",
            "Restored public trust through 100% transparent SLA resolution tracking."
        ]),
        (0.8, 3.1, 11.733, 1.7, "TIER 2: OPERATIONAL & TECHNICAL EFFICIENCY", LIGHT_GRAY, BLACK, [
            "Automated AI-driven triage reducing complaint routing latency from days to under 1 second.",
            "Exact GPS spatial coordinate mapping enabling precise municipal field inspector dispatch.",
            "Streamlined welfare scheme matching expanding government support accessibility."
        ]),
        (0.8, 5.0, 11.733, 1.7, "TIER 3: FUNCTIONAL CITIZEN BENEFITS", LIGHT_GRAY, BLACK, [
            "Simple photo + location complaint registration without complex administrative forms.",
            "Instant distress helpline access & real-time transportation hazard notifications.",
            "Gamified civic feedback upvoting empowering community-driven hazard prioritization."
        ])
    ]

    for tx, ty, tw, th, ttitle, tfill, tcolor, titems in tiers:
        box = slide15.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(tx), Inches(ty), Inches(tw), Inches(th))
        box.fill.solid()
        box.fill.fore_color.rgb = tfill
        box.line.color.rgb = NAVY_BLUE
        box.line.width = Pt(1.2)

        tf = box.text_frame
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.margin_top = Inches(0.15)
        tf.margin_left = Inches(0.3)
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = ttitle
        p.font.name = FONT_FAMILY
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = tcolor
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(4)

        for item in titems:
            p2 = tf.add_paragraph()
            p2.text = "• " + item
            p2.font.name = FONT_FAMILY
            p2.font.size = Pt(10.5)
            p2.font.color.rgb = tcolor
            p2.alignment = PP_ALIGN.LEFT

    # =========================================================================
    # SLIDE 16: 16. PERFORMANCE EVALUATION
    # =========================================================================
    slide16 = prs.slides.add_slide(blank_layout)
    add_header(slide16, "16. Performance Evaluation")

    pe_card = slide16.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.2), Inches(11.733), Inches(5.6))
    pe_card.fill.solid()
    pe_card.fill.fore_color.rgb = LIGHT_GRAY
    pe_card.line.color.rgb = BORDER_GRAY

    tf_pe = pe_card.text_frame
    tf_pe.vertical_anchor = MSO_ANCHOR.TOP
    tf_pe.margin_top = Inches(0.3)
    tf_pe.margin_left = Inches(0.4)
    tf_pe.margin_right = Inches(0.4)
    tf_pe.word_wrap = True

    p = tf_pe.paragraphs[0]
    p.text = "SUITABLE EVALUATION METRICS"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = NAVY_BLUE
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(14)

    pe_metrics = [
        ("AI / ML Metrics:", "Accuracy, Precision, Recall, F1-Score for automated complaint classification and department routing [TO BE MEASURED]."),
        ("Software Metrics:", "API Response Time [TO BE MEASURED], Page Load Speed [TO BE MEASURED], System Throughput & Reliability [TO BE MEASURED]."),
        ("Usability & Accessibility:", "Citizen task completion rate and multi-lingual user satisfaction score [TO BE MEASURED]."),
        ("Security & Auth Integrity:", "Verified zero-trust JWT authentication and password hashing security parameters.")
    ]

    for label, text in pe_metrics:
        p = tf_pe.add_paragraph()
        p.text = "• " + label + " "
        p.font.name = FONT_FAMILY
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = NAVY_BLUE
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = text
        run.font.bold = False
        run.font.color.rgb = BLACK
        p.space_after = Pt(12)

    # =========================================================================
    # SLIDE 17: 17. PROJECT TIMELINE
    # =========================================================================
    slide17 = prs.slides.add_slide(blank_layout)
    add_header(slide17, "17. Project Timeline")

    pt_box = slide17.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.2), Inches(11.733), Inches(5.6))
    pt_box.fill.solid()
    pt_box.fill.fore_color.rgb = LIGHT_GRAY
    pt_box.line.color.rgb = BORDER_GRAY

    tf_pt = pt_box.text_frame
    tf_pt.vertical_anchor = MSO_ANCHOR.TOP
    tf_pt.margin_top = Inches(0.3)
    tf_pt.margin_left = Inches(0.4)
    tf_pt.margin_right = Inches(0.4)
    tf_pt.word_wrap = True

    p = tf_pt.paragraphs[0]
    p.text = "PHASE I & PHASE II PROJECT TIMELINE"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = NAVY_BLUE
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(14)

    timeline_flow = [
        "Problem Identification ➔ Literature Survey ➔ Requirements Analysis ➔ System Architecture Design ➔ Development ➔ Testing & Integration ➔ Documentation",
        "Planned Milestones & Stages:",
        "• Phase 1 (Problem Identification & Literature Survey): Completed",
        "• Phase 2 (Requirement Analysis & Architecture Design): Completed",
        "• Phase 3 (Core Implementation & AI Integration): Completed",
        "• Phase 4 (System Testing & Performance Evaluation): In Progress",
        "• Phase 5 (Final Documentation & Review): Upcoming"
    ]

    for pt in timeline_flow:
        p = tf_pt.add_paragraph()
        p.text = pt
        p.font.name = FONT_FAMILY
        p.font.size = Pt(12)
        p.font.bold = True if "➔" in pt or "Planned" in pt else False
        p.font.color.rgb = NAVY_BLUE if "➔" in pt else BLACK
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(10)

    # =========================================================================
    # SLIDE 18: 18. WORK COMPLETED (Filled Boxes with Detailed Content - Image 4)
    # =========================================================================
    slide18 = prs.slides.add_slide(blank_layout)
    add_header(slide18, "18. Work Completed")

    milestones_rich = [
        (0.8, 1.2, "1. PROBLEM IDENTIFICATION", [
            "Identified fragmented paper, IVR, and multi-portal registration issues.",
            "Mapped administrative triage delays & lack of citizen SLA tracking.",
            "Defined single-window welfare scheme access requirements."
        ]),
        (4.8, 1.2, "2. LITERATURE SURVEY", [
            "Surveyed 5 IEEE research papers (2021-2024) on smart cities & NLP.",
            "Evaluated YOLOv5 pothole detection & municipal query chatbots.",
            "Analyzed e-governance complaint routing & rule-based engines."
        ]),
        (8.8, 1.2, "3. RESEARCH GAP", [
            "Identified missing links between reporting, AI triage & schemes.",
            "Found lack of regional Tamil/Tanglish multi-lingual support.",
            "Formulated CrowdCity AI multi-modal architecture."
        ]),
        (0.8, 4.1, "4. OBJECTIVES FORMULATION", [
            "Objective 1: Centralized citizen & authority ops platform.",
            "Objective 2: Groq AI automated complaint triage & NLP.",
            "Objective 3: GPS spatial mapping & multi-modal photo proof."
        ]),
        (4.8, 4.1, "5. METHODOLOGY & ARCHITECTURE", [
            "Formulated 11-step end-to-end process flowchart (Input to Tracking).",
            "Designed 6-layer system architecture (User, Backend, AI, Data).",
            "Established JWT security & Supabase PostgreSQL schema."
        ]),
        (8.8, 4.1, "6. STACK & PROTOTYPE MVP", [
            "Built Node.js & Express RESTful API MVC controllers.",
            "Integrated Groq LPU with openai/gpt-oss-120b & Vision AI.",
            "Developed Leaflet.js interactive map with real-time hazard pins."
        ])
    ]

    for mx, my, mtitle, mpts in milestones_rich:
        box = slide18.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(mx), Inches(my), Inches(3.733), Inches(2.6))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_GRAY
        box.line.color.rgb = NAVY_BLUE
        box.line.width = Pt(1.2)

        tf_m = box.text_frame
        tf_m.vertical_anchor = MSO_ANCHOR.TOP
        tf_m.margin_top = Inches(0.15)
        tf_m.margin_left = Inches(0.15)
        tf_m.margin_right = Inches(0.15)
        tf_m.word_wrap = True

        p = tf_m.paragraphs[0]
        p.text = mtitle
        p.font.name = FONT_FAMILY
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = NAVY_BLUE
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(6)

        for pt in mpts:
            p2 = tf_m.add_paragraph()
            p2.text = "• " + pt
            p2.font.name = FONT_FAMILY
            p2.font.size = Pt(9.5)
            p2.font.color.rgb = BLACK
            p2.alignment = PP_ALIGN.LEFT
            p2.space_after = Pt(3)

    # =========================================================================
    # SLIDE 19: 19. FUTURE WORK
    # =========================================================================
    slide19 = prs.slides.add_slide(blank_layout)
    add_header(slide19, "19. Future Work")

    roadmap_phases = [
        (0.8, "PHASE 1: NEAR-TERM REFINEMENT", [
            "Finalizing administrative authority dashboard controls.",
            "Optimizing multi-lingual voice recognition triggers.",
            "Refining rule-based scheme eligibility checker logic."
        ]),
        (4.8, "PHASE 2: PREDICTIVE ANALYTICS & VISION", [
            "Fine-tuning custom computer vision models for road damage.",
            "Integrating predictive analytics for urban decay trends.",
            "Adding flood, weather & traffic hotspot intelligence."
        ]),
        (8.8, "PHASE 3: PORTAL & MOBILE EXPANSION", [
            "Developing native Android & iOS mobile applications.",
            "Direct API integration with state e-Sevai portals.",
            "Conducting full municipal zone pilot deployments."
        ])
    ]

    for rx, rtitle, ritems in roadmap_phases:
        box = slide19.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(rx), Inches(1.2), Inches(3.733), Inches(5.6))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_GRAY
        box.line.color.rgb = NAVY_BLUE
        box.line.width = Pt(1.2)

        tf_r = box.text_frame
        tf_r.vertical_anchor = MSO_ANCHOR.TOP
        tf_r.margin_top = Inches(0.25)
        tf_r.margin_left = Inches(0.2)
        tf_r.margin_right = Inches(0.2)
        tf_r.word_wrap = True

        p = tf_r.paragraphs[0]
        p.text = rtitle
        p.font.name = FONT_FAMILY
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = NAVY_BLUE
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(12)

        for item in ritems:
            p2 = tf_r.add_paragraph()
            p2.text = "• " + item
            p2.font.name = FONT_FAMILY
            p2.font.size = Pt(11)
            p2.font.color.rgb = BLACK
            p2.alignment = PP_ALIGN.LEFT
            p2.space_after = Pt(10)

    # =========================================================================
    # SLIDE 20: 20. REFERENCES
    # =========================================================================
    slide20 = prs.slides.add_slide(blank_layout)
    add_header(slide20, "20. References")

    ref_box = slide20.shapes.add_textbox(Inches(0.8), Inches(1.15), Inches(11.733), Inches(5.8))
    tf20 = ref_box.text_frame
    tf20.word_wrap = True

    references = [
        "[1] S. Mellouer, A. El Amrani, and H. Bouzahir, \"Mobile crowdsensing and GIS integration for smart city hazard management,\" IEEE Access, vol. 11, pp. 45210-45222, 2023.",
        "[2] A. Kumar and R. Sharma, \"Automated municipal road hazard detection using deep convolutional neural networks,\" IEEE Transactions on Intelligent Transportation Systems, vol. 23, no. 8, pp. 12890-12901, 2022.",
        "[3] R. Sharma, P. Gupta, and V. Singh, \"Multilingual natural language processing for automated public grievance resolution in smart cities,\" in Proceedings of the IEEE International Conference on Smart Cities (ICSC), 2024, pp. 112-119.",
        "[4] M. Patel and S. Gupta, \"E-governance complaint classification and automated department routing using machine learning,\" IEEE Transactions on Engineering Management, vol. 70, no. 3, pp. 945-956, 2021.",
        "[5] K. Varma, R. Nair, and T. Sundaram, \"Automated welfare scheme eligibility processing using rule engines,\" IEEE Transactions on Computational Social Systems, vol. 10, no. 4, pp. 1820-1831, 2023.",
        "[6] Groq Cloud & Supabase Inc., \"Groq LPU Inference Engine and Supabase PostgreSQL Architecture Documentation,\" 2026. [Online]. Available: https://groq.com, https://supabase.com."
    ]

    for idx, ref in enumerate(references):
        p = tf20.paragraphs[0] if idx == 0 else tf20.add_paragraph()
        p.text = ref
        p.font.name = FONT_FAMILY
        p.font.size = Pt(12)
        p.font.color.rgb = BLACK
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(14)

    # =========================================================================
    # SLIDE 21: LITERATURE SURVEY – SUGGESTED FORMAT
    # =========================================================================
    slide21 = prs.slides.add_slide(blank_layout)
    add_header(slide21, "Literature Survey – Suggested Format")

    rows = 6
    cols = 6
    table_shape21 = slide21.shapes.add_table(rows, cols, Inches(0.5), Inches(1.15), Inches(12.333), Inches(5.75))
    t21 = table_shape21.table

    t21.columns[0].width = Inches(0.6)
    t21.columns[1].width = Inches(2.0)
    t21.columns[2].width = Inches(2.4)
    t21.columns[3].width = Inches(2.4)
    t21.columns[4].width = Inches(2.4)
    t21.columns[5].width = Inches(2.533)

    headers21 = ["S.No.", "Author & Year", "Paper / Topic", "Method / Technology", "Key Finding", "Limitation"]
    for idx, text in enumerate(headers21):
        cell = t21.cell(0, idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY_BLUE
        p = cell.text_frame.paragraphs[0]
        p.text = text
        p.font.name = FONT_FAMILY
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    lit_data21 = [
        ("1", "S. Mellouer et al.\n(2023)", "Mobile Crowdsensing & GIS Hazard Management", "Mobile Crowdsensing & GIS Mapping", "Enhanced spatial hazard mapping accuracy and citizen participation.", "High battery consumption; lacks automated AI triage."),
        ("2", "A. Kumar & R. Sharma\n(2022)", "Deep Learning for Municipal Road Hazard Detection", "CNN & Computer Vision (YOLOv5)", "Achieved 91% accuracy in detecting road craters from citizen photos.", "Limited to road hazards; lacks backend authority integration."),
        ("3", "R. Sharma et al.\n(2024)", "Multilingual NLP Chatbots in Public Admin", "LLMs & Retrieval-Augmented Generation", "Improved public query resolution response speed by 45%.", "English-only focus; lacks Tamil/Tanglish multi-lingual support."),
        ("4", "M. Patel & S. Gupta\n(2021)", "Automated Complaint Triage in E-Governance", "Text Classification & Random Forest", "Automated department routing for text-based citizen grievances.", "Inability to analyze multi-modal image evidence or GPS context."),
        ("5", "K. Varma et al.\n(2023)", "Integrated Welfare Scheme Eligibility Engines", "Rule-Based Engines & Indexing", "Streamlined welfare scheme matching for underprivileged citizens.", "Static rules without conversational AI guidance or tracking.")
    ]

    for row_idx, row_data in enumerate(lit_data21, start=1):
        for col_idx, cell_value in enumerate(row_data):
            cell = t21.cell(row_idx, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_GRAY if row_idx % 2 == 1 else WHITE
            p = cell.text_frame.paragraphs[0]
            p.text = cell_value
            p.font.name = FONT_FAMILY
            p.font.size = Pt(9.5)
            p.font.color.rgb = BLACK
            p.alignment = PP_ALIGN.CENTER if col_idx == 0 else PP_ALIGN.LEFT

    # =========================================================================
    # SLIDE 22: PROJECT TIMELINE – SUGGESTED FORMAT
    # =========================================================================
    slide22 = prs.slides.add_slide(blank_layout)
    add_header(slide22, "Project Timeline – Suggested Format")

    rows = 8
    cols = 4
    table_shape22 = slide22.shapes.add_table(rows, cols, Inches(0.8), Inches(1.2), Inches(11.733), Inches(5.6))
    t22 = table_shape22.table

    t22.columns[0].width = Inches(1.2)
    t22.columns[1].width = Inches(4.533)
    t22.columns[2].width = Inches(3.0)
    t22.columns[3].width = Inches(3.0)

    headers22 = ["Phase", "Activity", "Planned Date", "Status"]
    for idx, text in enumerate(headers22):
        cell = t22.cell(0, idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY_BLUE
        p = cell.text_frame.paragraphs[0]
        p.text = text
        p.font.name = FONT_FAMILY
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    timeline_template_rows = [
        ("1", "Problem Identification", "Aug 2026 – Sep 2026", "Completed"),
        ("2", "Literature Survey", "Sep 2026 – Oct 2026", "Completed"),
        ("3", "Requirement Analysis", "Oct 2026 – Nov 2026", "Completed"),
        ("4", "System Design", "Nov 2026 – Dec 2026", "Completed"),
        ("5", "Implementation", "Dec 2026 – Feb 2027", "Completed (MVP)"),
        ("6", "Testing & Evaluation", "Feb 2027 – Mar 2027", "In Progress"),
        ("7", "Documentation & Review", "Mar 2027 – Apr 2027", "Upcoming")
    ]

    for row_idx, (phase, activity, date_val, status_val) in enumerate(timeline_template_rows, start=1):
        c0 = t22.cell(row_idx, 0)
        c0.fill.solid()
        c0.fill.fore_color.rgb = LIGHT_GRAY if row_idx % 2 == 1 else WHITE
        p = c0.text_frame.paragraphs[0]
        p.text = phase
        p.font.name = FONT_FAMILY
        p.font.size = Pt(11)
        p.font.color.rgb = BLACK
        p.alignment = PP_ALIGN.CENTER

        c1 = t22.cell(row_idx, 1)
        c1.fill.solid()
        c1.fill.fore_color.rgb = LIGHT_GRAY if row_idx % 2 == 1 else WHITE
        p = c1.text_frame.paragraphs[0]
        p.text = activity
        p.font.name = FONT_FAMILY
        p.font.size = Pt(11)
        p.font.color.rgb = BLACK
        p.alignment = PP_ALIGN.LEFT

        c2 = t22.cell(row_idx, 2)
        c2.fill.solid()
        c2.fill.fore_color.rgb = LIGHT_GRAY if row_idx % 2 == 1 else WHITE
        p = c2.text_frame.paragraphs[0]
        p.text = date_val
        p.font.name = FONT_FAMILY
        p.font.size = Pt(11)
        p.font.color.rgb = BLACK
        p.alignment = PP_ALIGN.CENTER

        c3 = t22.cell(row_idx, 3)
        c3.fill.solid()
        c3.fill.fore_color.rgb = LIGHT_GRAY if row_idx % 2 == 1 else WHITE
        p = c3.text_frame.paragraphs[0]
        p.text = status_val
        p.font.name = FONT_FAMILY
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = NAVY_BLUE if "Completed" in status_val else BLACK
        p.alignment = PP_ALIGN.CENTER

    output_path = "CrowdCity_AI_Phase1_Review.pptx"
    try:
        prs.save(output_path)
        print(f"22-Slide Navy Blue Architectural presentation saved successfully to {os.path.abspath(output_path)}")
    except PermissionError:
        fallback_path = "CrowdCity_AI_Phase1_Review_Updated.pptx"
        prs.save(fallback_path)
        print(f"File {output_path} was locked. Saved successfully to {os.path.abspath(fallback_path)}")

if __name__ == "__main__":
    create_22_slide_diagram_deck_v2()
